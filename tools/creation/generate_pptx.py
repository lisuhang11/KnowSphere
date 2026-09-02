"""generate_pptx：根据结构化大纲生成 .pptx 并写入会话附件。"""

from __future__ import annotations

import io
import json
import logging
import re
from collections.abc import Mapping
from pathlib import Path
from typing import Annotated, Any

from langchain_core.runnables import RunnableConfig
from langchain_core.tools import InjectedToolArg, tool
from langgraph.prebuilt import ToolRuntime

from tools.events import emit_file_artifact, emit_thinking, emit_tool_call, emit_tool_result
from utils.run_config import thread_id_from_config

logger = logging.getLogger(__name__)

PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
MAX_SLIDES = 20
MAX_BULLETS = 8
MAX_TITLE_LEN = 80
MAX_BULLET_LEN = 200


def _clip(text: Any, limit: int) -> str:
    s = re.sub(r"\s+", " ", str(text or "")).strip()
    if len(s) <= limit:
        return s
    return s[: limit - 1].rstrip() + "…"


def _normalize_slides(raw: Any) -> list[dict[str, Any]]:
    items: list[Any]
    if isinstance(raw, str):
        try:
            parsed = json.loads(raw)
        except json.JSONDecodeError:
            parsed = []
        items = parsed if isinstance(parsed, list) else []
    elif isinstance(raw, list):
        items = raw
    else:
        items = []
    out: list[dict[str, Any]] = []
    for item in items[:MAX_SLIDES]:
        if isinstance(item, str):
            title = _clip(item, MAX_TITLE_LEN)
            bullets: list[str] = []
        elif isinstance(item, dict):
            title = _clip(item.get("title") or item.get("heading") or "", MAX_TITLE_LEN)
            raw_bullets = item.get("bullets") or item.get("points") or item.get("items") or []
            if isinstance(raw_bullets, str):
                raw_bullets = [raw_bullets]
            bullets = []
            if isinstance(raw_bullets, list):
                for b in raw_bullets[:MAX_BULLETS]:
                    line = _clip(b, MAX_BULLET_LEN)
                    if line:
                        bullets.append(line)
        else:
            continue
        if not title and not bullets:
            continue
        out.append({"title": title or "要点", "bullets": bullets})
    return out


def _safe_pptx_name(file_name: str, title: str) -> str:
    raw = (file_name or "").strip() or _clip(title, 40) or "演示文稿"
    name = Path(raw).name
    name = re.sub(r"[\\/:*?\"<>|]+", "_", name).strip(" .") or "演示文稿"
    if not name.lower().endswith(".pptx"):
        name = f"{name}.pptx"
    return name[:80]


def build_pptx_bytes(title: str, slides: list[dict[str, Any]]) -> bytes:
    from pptx import Presentation
    from pptx.util import Pt

    deck_title = _clip(title, MAX_TITLE_LEN) or "演示文稿"
    pages = _normalize_slides(slides)
    prs = Presentation()

    title_layout = prs.slide_layouts[0]
    cover = prs.slides.add_slide(title_layout)
    if cover.shapes.title is not None:
        cover.shapes.title.text = deck_title
    if len(cover.placeholders) > 1:
        cover.placeholders[1].text = "由 KnowSphere PPT 助手生成"

    body_layout = prs.slide_layouts[1]
    for page in pages:
        slide = prs.slides.add_slide(body_layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = page["title"]
        try:
            body = slide.placeholders[1]
        except (KeyError, IndexError):
            continue
        tf = body.text_frame
        tf.clear()
        bullets = page["bullets"] or ["（无要点）"]
        for i, line in enumerate(bullets):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.text = line
            para.level = 0
            para.font.size = Pt(18)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _config_from_runtime(runtime: ToolRuntime | None) -> dict[str, Any] | None:
    if runtime is None:
        return None
    cfg = getattr(runtime, "config", None)
    if isinstance(cfg, dict):
        return cfg
    if isinstance(cfg, Mapping):
        return dict(cfg)
    return None


def _session_id(config: RunnableConfig | None, runtime: ToolRuntime | None) -> str | None:
    return thread_id_from_config(config) or thread_id_from_config(_config_from_runtime(runtime))


def _save_pptx(session_id: str, file_name: str, data: bytes) -> dict[str, Any]:
    from utils.temporary_attachments import TemporaryAttachmentStore

    store = TemporaryAttachmentStore()
    row = store.create(
        session_id=session_id,
        file_name=file_name,
        mime_type=PPTX_MIME,
        file_size=len(data),
        data=data,
    )
    try:
        store.mark_ready(row["id"], content=f"[generated pptx] {file_name}")
    except Exception:
        logger.debug("标记生成 PPT 为 ready 失败", exc_info=True)
    return row


@tool
def generate_pptx(
    title: str,
    slides: list[dict],
    file_name: str = "",
    runtime: Annotated[ToolRuntime | None, InjectedToolArg] = None,
    config: Annotated[RunnableConfig | None, InjectedToolArg] = None,
) -> str:
    """根据标题和每页要点生成 PPTX 文件。材料准备好后再调用。

    slides 为对象列表，每项含 title（页标题）和 bullets（要点字符串列表）。
    不要一次塞入超过 20 页。用户要修改时重新调用本工具生成完整新文件。
    """
    writer = getattr(runtime, "stream_writer", None) if runtime is not None else None
    deck_title = _clip(title, MAX_TITLE_LEN) or "演示文稿"
    pages = _normalize_slides(slides)
    emit_tool_call("generate_pptx", f"正在生成「{deck_title}」…", writer)
    emit_thinking(f"【生成 PPT】{deck_title}，{len(pages)} 页内容", writer)

    if not pages:
        msg = "没有可用的幻灯片内容，请提供每页 title 和 bullets。"
        emit_tool_result("generate_pptx", msg, success=False, writer=writer)
        return json.dumps({"ok": False, "message": msg}, ensure_ascii=False)

    session_id = _session_id(config, runtime)
    if not session_id:
        msg = "无法保存文件：缺少会话。"
        emit_tool_result("generate_pptx", msg, success=False, writer=writer)
        return json.dumps({"ok": False, "message": msg}, ensure_ascii=False)

    try:
        data = build_pptx_bytes(deck_title, pages)
        saved_name = _safe_pptx_name(file_name, deck_title)
        row = _save_pptx(session_id, saved_name, data)
    except Exception as exc:
        logger.exception("生成 PPT 失败")
        msg = f"生成 PPT 失败：{exc}"
        emit_tool_result("generate_pptx", msg, success=False, writer=writer)
        return json.dumps({"ok": False, "message": msg}, ensure_ascii=False)

    artifact = {
        "id": row["id"],
        "file_name": row["file_name"],
        "file_type": row.get("file_type") or "pptx",
        "file_size": int(row.get("file_size") or len(data)),
        "mime_type": row.get("mime_type") or PPTX_MIME,
    }
    emit_file_artifact(
        attachment_id=artifact["id"],
        file_name=artifact["file_name"],
        file_type=artifact["file_type"],
        file_size=artifact["file_size"],
        mime_type=artifact["mime_type"],
        writer=writer,
    )
    summary = f"已生成 {artifact['file_name']}，共 {len(pages) + 1} 页（含封面）。"
    emit_tool_result("generate_pptx", summary, writer=writer)
    return json.dumps(
        {"ok": True, "message": summary, "artifact": artifact, "slide_count": len(pages) + 1},
        ensure_ascii=False,
    )
