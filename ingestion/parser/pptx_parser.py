"""pptx 链式解析器（docx2：MarkItDown 优先链，Q7）。

- PptxParser：python-pptx 专用解析器（逐页提取文本框/表格/图片）。
- Pptx2Parser：FirstParser([MarkItDownParser, PptxParser])，与 Docx2Parser 同构。
"""

from __future__ import annotations

import logging

from ingestion.parser.base_parser import BaseParser, ParserError, ParseResult
from ingestion.parser.chain_parser import FirstParser
from ingestion.parser.markitdown_parser import MarkItDownParser

logger = logging.getLogger(__name__)


class PptxParser(BaseParser):
    supported_file_types = ["pptx"]

    def parse(self, path: str) -> ParseResult:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise ParserError("python-pptx 未安装，请执行 pip install python-pptx") from exc

        try:
            prs = Presentation(path)
        except Exception as exc:  # noqa: BLE001
            raise ParserError(f"python-pptx 打开失败: {exc}") from exc

        result = ParseResult()
        parts: list[str] = []
        for idx, slide in enumerate(prs.slides, start=1):
            slide_parts: list[str] = []
            try:
                for shape in slide.shapes:
                    self._handle_shape(shape, slide_parts, result)
            except Exception as exc:  # noqa: BLE001
                logger.debug("pptx slide %d failed: %s", idx, exc)
            if slide_parts:
                parts.append(f"## 幻灯片 {idx}\n\n" + "\n\n".join(slide_parts))
        result.markdown = "\n\n".join(parts).strip()
        if not result.markdown:
            raise ParserError("pptx 未解析出任何文本")
        return result

    def _handle_shape(self, shape, parts: list[str], result: ParseResult) -> None:
        try:
            if shape.has_text_frame:
                text = "\n".join(p.text for p in shape.text_frame.paragraphs if p.text)
                if text.strip():
                    parts.append(text.strip())
            if shape.shape_type == 13:  # PICTURE
                image = shape.image
                data = image.blob
                mime = getattr(image, "content_type", "image/png") or "image/png"
                parts.append(self._embed_image(result, data, mime, original_ref=getattr(shape, "name", "")))
            if shape.has_table:
                rows: list[str] = []
                for i, row in enumerate(shape.table.rows):
                    cells = [" ".join(c.text.split()) for c in row.cells]
                    rows.append("| " + " | ".join(cells) + " |")
                    if i == 0:
                        rows.append("| " + " | ".join("---" for _ in cells) + " |")
                if rows:
                    parts.append("\n".join(rows))
            if shape.shape_type == 6:  # GROUP
                for sub in shape.shapes:
                    self._handle_shape(sub, parts, result)
        except Exception as exc:  # noqa: BLE001
            logger.debug("pptx shape %s failed: %s", getattr(shape, "shape_id", "?"), exc)


class Pptx2Parser(FirstParser):
    def __init__(self, parse_options: dict | None = None, *, markitdown_first: bool = True):
        if parse_options:
            self.parse_options = parse_options
        md_parser = MarkItDownParser(parse_options)
        pptx_parser = PptxParser(parse_options)
        parsers = [md_parser, pptx_parser] if markitdown_first else [pptx_parser, md_parser]
        super().__init__(parsers, parse_options)


__all__ = ["PptxParser", "Pptx2Parser"]
